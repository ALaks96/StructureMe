# Script containing all the programs to extract text, images, tables from a variety of file types
import re
import os
import cv2
import zipfile
import pytesseract
import numpy as np
import pandas as pd
from PIL import Image
try:
    from xml.etree.cElementTree import XML
except ImportError:
    from xml.etree.ElementTree import XML
from io import StringIO
from pptx import Presentation
from pdfminer.pdfpage import PDFPage
from pdfminer.layout import LAParams
from pdf2image import convert_from_path
from Preprocessing.beautify import fix_text
from pdfminer.converter import TextConverter
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
WORD_NAMESPACE = '{http://schemas.openxmlformats.org/wordprocessingml/2006/main}'
PARA = WORD_NAMESPACE + 'p'
TEXT = WORD_NAMESPACE + 't'


def docxExtractor(path):
    # Function to extract content from docx files, takes path input if path endswith .docx
    # Start by reading MSoffice zipfile
    document = zipfile.ZipFile(path)
    # Search in xml structure the location of content (here it's word/document.xml)
    xml_content = document.read('word/document.xml')
    document.close()
    # Generate xml tree structure from content location
    tree = XML(xml_content)
    # Initialize dictionary to contain content and index per paragraph
    doc = {}
    # Initialize string to contain concatenated text from document for nlp purposes
    s = ''
    # vector = {}
    paragraph_nb = 1
    # Iterate through all elements of the xml tree
    for paragraph in tree.getiterator(PARA):
        # Append to list if node in tree contains non-null text
        texts = [node.text
                 for node in paragraph.getiterator(TEXT)
                 if node.text]
        if texts:
            # Concatenate non null text contained in previous list
            text = ''.join(texts)
            # Index concatenated string to paragtaph number
            doc[str(paragraph_nb)] = fix_text(text)
            # Append concatenated string to current string (for nlp)
            s += fix_text(text)
#            if vectors:
#                vector[str(paragraph_nb)] = vectorizer(text, lang=detect(text))
#            else:
#                pass
            paragraph_nb += 1

#    if vectors:
#        return creator, doc, vector
#    else:
    return doc, s


def ppt_extractor(path):
    # Initialize dictionary to contain content of pptx per slide
    paragraph_repo = {}
    # vector = {}
    f = open(path, "rb")
    # Use Presentation module from python-pptx to process content
    prs = Presentation(f)
    slide_nb = 0
    # Initialize string to contain concatenated text from document for nlp purposes
    s = ''
    for slide in prs.slides:

        slide_nb += 1
        # Initialize temporary text container
        temp_text = ''

        for shape in slide.shapes:
            # If string is not null and is text, append to temporary text container
            if hasattr(shape, "text") and shape.text.strip():
                temp_text += shape.text
        # Once all text has been retrieved from slide attribute it to page nb index inside initialized dictionary
        paragraph_repo[str(slide_nb)] = fix_text(temp_text)
        # append concatenated string to initialized string for nlp purposes
        s += fix_text(temp_text)
        # if vectors:
        #     vector[str(slide_nb)] = vectorizer(temp_text, lang=detect(text))
        # else:
        #     pass

    # if vectors:
    #     return creator, paragraph_repo, vector
    # else:
    return paragraph_repo, s


def txt_extractor(path):
    # Initialize dictionary to contain all content of txt file
    doc = {}
    # vector = {}
    paragraph_nb = 1
    # Initialize string to contain all concatenated text from txt file
    s = ""

    with open(path) as f:
        lines = f.read()

    # Split document at line jumps
    texts = lines.strip().split("/n/n")
    for text in texts:
        # Append text to line index
        doc[str(paragraph_nb)] = fix_text(text)
        # Append concatenated string to intialized string for nlp purposes
        s += fix_text(text)
        # if vectors:
        #     vector[str(paragraph_nb)] = vectorizer(text, lang=detect(text))
        # else:
        paragraph_nb += 1

    # if vectors:
    #     return doc, vector
    # else:
    return doc, s


def pdf_extractor(path):
    # Initialize PDFResourceManager to go through content of PDF
    rsrcmgr = PDFResourceManager()
    retstr = StringIO()
    laparams = LAParams()
    # Initialize TextConverter method to convert PDF data to text
    device = TextConverter(rsrcmgr, retstr, laparams=laparams)
    fp = open(path, 'rb')
    interpreter = PDFPageInterpreter(rsrcmgr, device)
    # Password may be required to open PDF
    password = ""
    maxpages = 0
    caching = True
    pagenos = set()
    current_page_number = 1
    # Initialize dictionary to contain content of PDF indexed per page
    paragraph_repo = {}
    # vector = {}
    # Initialize string to contain all concatenated text for further nlp processing
    s = ''

    # Loop through every page the PDFResourceManager has identified with listed options above
    for page in PDFPage.get_pages(fp, pagenos, maxpages=maxpages, password=password, caching=caching,
                                  check_extractable=True):
        interpreter.process_page(page)
        # Get text value
        text = retstr.getvalue()
        retstr.truncate(0)
        # Fix potential encoding error
        text = re.sub(u'(\u0000)', "", text)
        # Attribute retrieved text to current page number in intialized dictionary
        paragraph_repo[str(current_page_number)] = fix_text(text)
        # Append concatenated string to initialized string for nlp
        s += fix_text(text)
        # if vectors:
        #     vector[str(current_page_number)] = vectorizer(text, lang=detect(text))
        # else:
        #     pass
        current_page_number += 1

    fp.close()
    device.close()
    retstr.close()
    # if vectors:
    #     return Classified, creator, paragraph_repo, vector
    # else:
    return paragraph_repo, s


def img_extractor(path):
    # initialize dictionary to contain image contents
    dic = {}
    # Read image
    img = cv2.imread(path)
    # Attempt to convert image to string in case of a scanned document non registred as a PDF
    scan = pytesseract.image_to_string(img)
    # Text obtained contains a certain amount of characters, consider it a scan
    if scan or len(scan) > 20:
        dic["scanned_document"] = scan
    # Otherwise consider it as a photo from which we can not extract its content (look out for summarization script)
    else:
        dic["photo"] = "Content of this photo to be classified"
    # Return dictionnary and raw img content for Object Detection purposes
    return dic, img


def excel_extractor(path):
    # Initialize dictionary to contain content of every potential sheet within an excel spreadsheet
    dic = {}
    filename = os.path.basename(path)
    # Verify it is indeed an excel file
    if filename.endswith(".xlsx") or filename.endswith(".xls"):
        table = pd.read_excel(path, sheet_name=None)
        # For every sheet, replace NaN value by null to avoid parsing errors in JSON format
        # Rearrange content of sheets in record-style dictionaries
        for k in table.keys():
            table[k] = table[k].replace({np.nan: None})
            dic[k] = table[k].to_dict(orient='records')
    # Return dictionary and raw table for table characteristics extraction
        return dic, table
    else:
        return table_extractor(path)


def table_extractor(path):
    # Initialize dictionary to contain contents of any type of table except excel
    dic = {}
    filename = os.path.basename(path)
    # For csv's and tsv's use pandas method
    if filename.endswith("csv") or filename.endswith("tsv"):
        table = pd.read_csv(path, encoding='utf-8')
        # Fix a posteriori index_col initialization problem when using pandas method
        if table.columns[0] == 'Unnamed: 0':
            del table['Unnamed: 0']
        # replace NaN value by null to avoid parsing errors in JSON format
        table = table.replace({np.nan: None})
        # Rearrange content of table in record-style dictionaries
        dic["single_table"] = table.to_dict(orient='records')
    # For other table formats use slower more general purpose pandas method
    else:
        table = pd.read_table(path)
        # Same process as above
        if table.columns[0] == 'Unnamed: 0':
            del table['Unnamed: 0']
        table = table.replace({np.nan: None})
        dic["single_table"] = table
    # Return dictionary and raw table for table characteristics extraction
    return dic, table


def scan_extractor(path):
    # Initialize dictionary to contain contents of scanned document that PDFResourceManager failed to handle
    paragraph_repo = {}
    # Initialize string to contain concatenated string for nlp purposes
    s = ''
    # Store all the pages of the PDF in a variable
    pages = convert_from_path(path, 500)

    # Counter to store images of each page of PDF to image
    image_counter = 1
    photos = []
    # Iterate through all the pages stored above
    for page in pages:
        # Declaring filename for each page of PDF as JPG
        # For each page, filename will be:
        # PDF page 1 -> page_1.jpg
        # PDF page 2 -> page_2.jpg
        # PDF page 3 -> page_3.jpg
        # ....
        # PDF page n -> page_n.jpg
        filename = "page_" + str(image_counter) + ".jpg"
        photos.append(filename)

        # Save the image of the page in system
        page.save(filename, 'JPEG')

        # Increment the counter to update filename
        image_counter = image_counter + 1

    ''' 
    Part #2 - Recognizing text from the images using OCR 
    '''
    # Variable to get count of total number of pages
    filelimit = image_counter - 1

    # Iterate from 1 to total number of pages
    for i in range(1, filelimit + 1):
        # Set filename to recognize text from
        # Again, these files will be:
        # page_1.jpg
        # page_2.jpg
        # ....
        # page_n.jpg
        filename = "page_" + str(i) + ".jpg"

        # Recognize the text as string in image using pytesserct
        text = str((pytesseract.image_to_string(Image.open(filename))))

        # The recognized text is stored in variable text
        # Any string processing may be applied on text
        # Here, basic formatting has been done:
        # In many PDFs, at line ending, if a word can't
        # be written fully, a 'hyphen' is added.
        # The rest of the word is written in the next line
        # Eg: This is a sample text this word here GeeksF-
        # orGeeks is half on first line, remaining on next.
        # To remove this, we replace every '-\n' to ''.
        text = fix_text(text.replace('-\n', ''))
        paragraph_repo[str(i)] = text
        s += text

    for file in photos:
        os.remove(file)

    return paragraph_repo, s


def get_content(path):
    """
    General purpose function to handle content extraction based on file type.
    Will call one of the above functions.
    """
    text = {}
    global raw, file_type
    if path.endswith(".pptx") or path.endswith(".ppt"):
        # try:
        text, raw = ppt_extractor(path)
        file_type = "txt"
        # except Exception as e:
        #     print(e)
        #     pass

    elif path.endswith(".pdf"):

        # try:
        text, raw = pdf_extractor(path)
        if raw == '':
            text, raw = scan_extractor(path)
        file_type = "txt"
        # except Exception as e:
        #     print(e)
        #     try:
        #         text = scan_extractor(path)
        #     except Exception as ee:
        #         print(ee)
        #         pass

    elif path.endswith(".docx"):

        # try:
        text, raw = docxExtractor(path)
        file_type = "txt"
        # except Exception as e:
        #     print(e)
        #     pass

    elif path.endswith(".txt"):

        # try:
        text, raw = txt_extractor(path)
        file_type = "txt"
        # except Exception as e:
        #     print(e)
        #     pass

    elif path.endswith(".png") \
            or path.endswith(".jpeg") \
            or path.endswith(".jpg"):

        text, raw = img_extractor(path)
        file_type = "img"

    elif path.endswith(".tsv") \
            or path.endswith(".csv"):
        # try:
        text, raw = table_extractor(path)
        file_type = "table"
        # except Exception as e:
        #     print(e)

    elif path.endswith(".xls")\
            or path.endswith(".xlsx"):

        text, raw = excel_extractor(path)
        file_type = "sheets"

    else:
        pass
    return text, raw, file_type
